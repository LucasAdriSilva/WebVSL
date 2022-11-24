from flask import Flask, render_template, flash, redirect, request,  url_for, session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from waitress import serve
from flask_pymongo import PyMongo
import shutil
import envio
import os

app = Flask(__name__)

app.config["MONGO_URI"] = "mongodb+srv://Lucas:Uche3517@cluster0.xghdh47.mongodb.net/CVS"
app.config["SECRET_KEY"]= "hello"
mongo = PyMongo(app)


@app.route("/")
def index():
    return render_template("principal.html")

@app.route("/del")
def deletField():
    name = session["username"]
    dir = f"static/{name}"
    shutil.rmtree(dir)
    return render_template("principal.html")

@app.route("/convert")
def conversor():
  if "username" in session:
    return render_template("conversor.html")
  else:
    return render_template("login.html")

@app.route('/autenticar', methods=['GET'])
def autenticar():
  if "username" in session:
    vsl = request.args.get('vsl')
    vsl = vsl.split('\n')
    email = request.args.get('email')
    email = email.split(',')
    x = []

    #Faz o tratamento da string enviada pelo front
    #Add em um array e se estiver espaço vazio é excluido automaticamente
    for item in vsl:
      st = item.replace('\r', '')
      st = st.strip()
      x.append(st)
    
    #Remove linhas vazias
    while ("" in x):
      x.remove("")

    try:
      #Abrindo PowerPoint
      apresentacao = Presentation()
      apresentacao.slide_width = Inches(16)
      apresentacao.slide_height = Inches(9)

      # Cada item no array ira fazer esse bloco
      for line in x:

        #Criando um slide
        slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

        #Criando e centralizando a TextBox
        x = Inches(3)
        y = Inches(3.5)
        largura = Inches(10)
        altura = Inches(2)
        caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)
        caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

        #Formatando o TextBox
        caixa_texto.text = line
        caixa_texto.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        caixa_texto.text_frame.paragraphs[0].font.size = Pt(40)
        caixa_texto.text_frame.paragraphs[0].font.name = "Arial"
        caixa_texto.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        caixa_texto.text_frame.word_wrap = True

      #Salvando o arquivo pptx
      apresentacao.save("VSL.pptx")

      #Salvando em um diretorio proprio da sessão
      name = session["username"]
      dir = f"./static/{name}/"
      os.makedirs(dir)


      #Alterando o arquivo para um diretorio unico da sessão
      diretorio = f"static/{name}/VSL.pptx"
      shutil.move('VSL.pptx', diretorio)

      if email != None:
        envio.send_email_with_anexo(
          email,
          "Envio automatico da VSL, disponviel para download em anexo desse e-mail",
          "Envio da VSL gerada",
          diretorio
        )
    except Exception as e:
      print(e)
  else:
    return render_template("login.html")

  return render_template("resultado.html", dir = f"static/{name}/VSL.pptx")

@app.route('/reset', methods=["GET","POST"])
def reset():
  try:  
    if request.method == "POST":
      login = request.form.get("login")
      user_found = mongo.db.user.find_one({"name": login})

      if user_found:
        newPassword = envio.random_generator()
        mongo.db.user.find_one_and_update({'_id': user_found["_id"]}, {"$set": {'password': newPassword}})

        content = f"Segue a sua nova senha {newPassword}"
        subject = "Alteração de senha"
        envio.send_email(user_found["email"], content, subject)

        return render_template("login.html", msgTrue ="E-mail com a nova senha enviado")
      else:
        return render_template("reset.html", msgFalse ="Login não encontrado")
    else:
      return render_template("reset.html")
  except  Exception as e:
    print(e)

@app.route('/togglePassword', methods=["GET","POST"])
def togglePassword():
  try:
    if "username" in session:
      if request.method == "POST":
        oldPassword = request.form.get("oldPassword")
        newPassword = request.form.get("newPassword")
        passwordConfirmation = request.form.get("passwordConfirmation")
        
        # if oldPassword == user.password:
        user = mongo.db.user.find_one({"name": session['username']})
        # else:
        #   return render_template("togglePassword.html", msgFalse= 'Usuário não encontrado')
        print (f'{newPassword}, {passwordConfirmation}')
        if newPassword == passwordConfirmation:
          mongo.db.user.find_one_and_update({'_id': user["_id"]}, {"$set": {'password': newPassword}})
          return render_template("login.html", msg= 'Senha alterada com sucesso')
        else:
          return render_template("togglePassword.html", msgFalse= 'As senhas não são) iguais, por favor verificar')
        if user.password == passs:
          return "Mesma senha, comfirmada"
        else:
          print("Deu ruim")
      else:
        return render_template("togglePassword.html")
    else:
      return render_template("login.html")
  except  Exception as e:
    print(e)

  except  Exception as e:
    print(e)

@app.route('/register', methods=["GET", "POST"])
def register():
  res_valided = "Usuário cadastrado"
  if request.method == "POST":
    username = request.form.get("usuario")
    email = request.form.get("email")
    password = request.form.get("password")

    user_found = mongo.db.user.find_one({"name": username})

    if user_found:
      flash("O usuário já existe", "error")
      return render_template("register.html")
    else:
      mongo.db.user.insert_one({
        "name": str(username),
        "password": str(password),
        "email": str(email)
      })
      return render_template("login.html", msg = res_valided)
  else:
    return render_template("register.html")

@app.route("/login", methods=["GET","POST"])
def login():
  try:
    #verifica se está logado
    if "username" in session:
        # return redirect(url_for("principal.html"))
        return render_template("principal.html")

    #Se for um login executara esse trecho
    elif request.method == "POST":
        username = request.form.get("usuario")
        password = request.form.get("senha")

        user_found = mongo.db.user.find_one({"name": username})

        if user_found:
            user_val = user_found["name"]
            passwordcheck = user_found["password"]

            # if check_password_hash(passwordcheck, password):
            if passwordcheck == password:
                session["username"] = user_val
                print(session["username"])
                # return redirect(url_for("principal.hmtl"))
                return render_template("conversor.html")
            else:
                flash("Senha Incorreta", "error")
                return render_template("login.html")
        else:
            flash("Usuário não encontrado")
            return render_template("login.html")

    elif request.method == "GET":
      return render_template("login.html")

  except Exception as e:
        print(f"retorno é {e}")

@app.route("/logout")
def logout():
    session.pop("username", None)
    return redirect(url_for("login"))

if __name__ == "__main__":
  app.run(debug=True)
  # serve(app, host="0.0.0.0", port=8080)
